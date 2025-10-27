import streamlit as st
from PIL import Image
from pptx import Presentation
import tempfile


def publish_vs_post():

    with st.expander("Template path"):
        template_path = st.file_uploader(
            key="publish_vs__template_path",
            label="Choose a template",
            type=["pptx"]
        )

        if template_path is not None and st.button("🛠️ Générer la présentation"):
            # Sauvegarder le fichier uploadé dans un fichier temporaire
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
                tmp_file.write(template_path.getvalue())
                template_path = tmp_file.name

    with st.expander("Home logo"):
        home_logo = st.file_uploader(
            key="publish_vs__logo__home",
            label="Choose a logo...",
            type=["jpg", "jpeg", "png"]
        )

        if home_logo:
            home_image = Image.open(home_logo)
        else:
            home_image = None

    with st.expander("Away logo"):
        away_logo = st.file_uploader(
            key="publish_vs__logo__away",
            label="Choose a logo...",
            type=["jpg", "jpeg", "png"]
        )

        if away_logo:
            away_image = Image.open(away_logo)
        else:
            away_image = None

    with st.expander("First slide"):
        team_home = st.text_input(
            key="publish_vs__team__home",
            label="Home Team"
        )

        team_away = st.text_input(
            key="publish_vs__team__away",
            label="Away Team"
        )

    with st.expander("Away slides"):
        nb_slides_away = st.number_input(
            key="publish_vs__nb_slides__away",
            label="How many slides for the away team?",
            min_value=0,
            max_value=10,
            value=1,
            step=1
        )

        slides_text_away = []
        for i in range(nb_slides_away):
            with st.expander(f"Content of the away slide {i + 1}"):
                title = st.text_input("Title", key=f"publish_vs__away_slide_{i}__title")
                stat = st.text_input("Stat", key=f"publish_vs__away_slide_{i}__stat")
                sub_title = st.text_input("Sub title", key=f"publish_vs__away_slide_{i}__sub_title")
                sub_sub_title = st.text_input("Sub sub title", key=f"publish_vs__away_slide_{i}__sub_sub_title")

                slides_text_away.append({
                    "title": title,
                    "stat": stat,
                    "stat_shape": stat,
                    "sub_title": sub_title,
                    "sub_sub_title": sub_sub_title
                })

    with st.expander("Intermediate slide"):
        with st.expander("Home team"):
            text_home = st.text_input("Text", key="publish_vs__intermediate__home")
        with st.expander("Away team"):
            text_away = st.text_input("Text", key="publish_vs__intermediate__away")

        intermediate = {
            "home": text_home,
            "away": text_away
        }

    with st.expander("Home slides"):
        nb_slides_home = st.number_input(
            key="publish_vs__nb_slides__home",
            label="How many slides for the home team?",
            min_value=0,
            max_value=10,
            value=1,
            step=1
        )

        slides_text_home = []
        for i in range(nb_slides_home):
            with st.expander(f"Content of the home slide {i + 1}"):
                title = st.text_input("Title", key=f"publish_vs__home_slide_{i}__title")
                stat = st.text_input("Stat", key=f"publish_vs__home_slide_{i}__stat")
                sub_title = st.text_input("Sub title", key=f"publish_vs__home_slide_{i}__sub_title")
                sub_sub_title = st.text_input("Sub sub title", key=f"publish_vs__home_slide_{i}__sub_sub_title")

                slides_text_home.append({
                    "title": title,
                    "stat": stat,
                    "stat_shape": stat,
                    "sub_title": sub_title,
                    "sub_sub_title": sub_sub_title
                })

    if st.button(
        key="publish_vs__button",
        label="🚀 Générer le PowerPoint et les images",
        type="primary"
    ):
        # st.image(home_image)
        # st.write(
        #     team_away,
        #     slides_text_away
        # )

        # st.write(intermediate)

        # st.image(away_image)
        # st.write(
        #     team_away,
        #     slides_text_away
        # )

        generate_pptx_from_template(
            template_path,
            "",
            team_home,
            team_away,
            home_image,
            away_image,
            slides_text_home,
            slides_text_away,
            intermediate
        )

        # st.session_state.powerpoint_path, st.session_state.slides_png_paths = generer_ppt_et_png(
        #     template_path, textes_slides
        # )

    # # --- Étape 4: Publication ---
    # if st.session_state.slides_png_paths:
    #     st.header("4. Publier sur Instagram")
    #     caption = st.text_area("Légende pour le post Instagram", height=100, value="#football #stats")

    #     if st.button("📤 Publier sur Instagram"):
    #         st.info("Connexion à Instagram et publication...")
    #         # Appel de la fonction de publication (voir Phase 4)
    #         publier_sur_instagram(st.session_state.slides_png_paths, caption)
    #         st.success("Post publié avec succès !")


def generate_pptx_from_template(
        template_path,
        output_path,
        team_home,
        team_away,
        home_image,
        away_image,
        slides_text_home,
        slides_text_away,
        intermediate
):
    prs = Presentation(template_path)

    for i, slide in enumerate(prs.slides):
        st.write(i, slide)
        # if i < len(textes_slides):
        #     data = textes_slides[i]  # dictionnaire : {"titre": "...", "contenu": "..."}

        #     # Parcourir les formes de la diapositive
        #     for shape in slide.shapes:
        #         if not shape.has_text_frame:
        #             continue

        #         # Exemple : remplacement selon un mot-clé dans le texte existant
        #         if "{titre}" in shape.text:
        #             shape.text = data.get("titre", "")
        #         elif "{contenu}" in shape.text:
        #             shape.text = data.get("contenu", "")

    # prs.save(output_path)
    # return output_path
